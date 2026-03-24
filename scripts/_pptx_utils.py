#!/usr/bin/env python3
from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation


def load_prs(path: str | Path) -> Presentation:
    return Presentation(str(path))


def shape_target(shape, shape_index: int) -> dict[str, Any]:
    target = {
        "shape_index": shape_index,
        "shape_name": shape.name,
    }
    if shape.is_placeholder:
        pf = shape.placeholder_format
        target["placeholder_idx"] = pf.idx
        target["placeholder_type"] = str(pf.type)
    return target


def paragraph_to_dict(paragraph) -> dict[str, Any]:
    text = paragraph.text or ""
    item: dict[str, Any] = {"text": text}
    if paragraph.level:
        item["level"] = paragraph.level
    if paragraph.runs:
        bold_values = [run.font.bold for run in paragraph.runs if run.text]
        if bold_values and all(v is True for v in bold_values):
            item["bold"] = True
        elif bold_values and all(v is False for v in bold_values):
            item["bold"] = False
    return item


def shape_text_payload(shape) -> dict[str, Any] | None:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return None
    paragraphs = [paragraph_to_dict(p) for p in shape.text_frame.paragraphs]
    return {
        "paragraphs": paragraphs,
        "text": "\n".join(p["text"] for p in paragraphs),
    }


def shape_image_payload(shape) -> dict[str, Any] | None:
    image = getattr(shape, "image", None)
    if image is None:
        return None
    return {
        "image_filename": image.filename,
        "image_ext": image.ext,
        "image_size": list(image.size),
        "image_sha1": image.sha1,
    }


def choose_shape(slide, target: dict[str, Any]):
    if "shape_name" in target:
        for shape in slide.shapes:
            if shape.name == target["shape_name"]:
                return shape
    if "shape_index" in target:
        idx = target["shape_index"]
        if 0 <= idx < len(slide.shapes):
            return slide.shapes[idx]
    if "placeholder_idx" in target:
        wanted = target["placeholder_idx"]
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == wanted:
                return shape
    raise KeyError(f"Cannot find shape for target {target}")


def replace_shape_image(slide, shape, image_path: str | Path) -> None:
    if getattr(shape, "image", None) is None:
        raise TypeError(f"Shape does not contain an image: {shape.name}")
    _, rel_id = slide.part.get_or_add_image_part(str(image_path))
    shape._pic.blipFill.blip.rEmbed = rel_id


def write_payload(path: str | Path, payload: Any) -> None:
    out = Path(path)
    if out.suffix.lower() == ".json":
        out.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n")
    else:
        out.write_text(yaml.safe_dump(payload, allow_unicode=True, sort_keys=False), encoding="utf-8")


def read_payload(path: str | Path) -> Any:
    p = Path(path)
    text = p.read_text(encoding="utf-8")
    if p.suffix.lower() == ".json":
        return json.loads(text)
    return yaml.safe_load(text)
